#!/usr/bin/env python3
import sys
import os
import argparse
import io

# Prefer PyMuPDF (fitz) for extracting text and images from PDF pages
try:
    import fitz  # PyMuPDF
except Exception:
    fitz = None

try:
    from docx import Document as DocxDocument
    from docx.shared import Inches as DocxInches
except Exception:
    DocxDocument = None
    DocxInches = None

try:
    from openpyxl import Workbook
    from openpyxl.drawing.image import Image as XLImage
except Exception:
    Workbook = None
    XLImage = None

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except Exception:
    Presentation = None


def extract_with_fitz(input_path):
    if fitz is None:
        return None
    doc = fitz.open(input_path)
    pages = []
    for p in doc:
        page_text = p.get_text("text")
        images = []
        for img_index, img in enumerate(p.get_images(full=True)):
            xref = img[0]
            base_image = doc.extract_image(xref)
            image_bytes = base_image["image"]
            image_ext = base_image.get("ext", "png")
            images.append({"bytes": image_bytes, "ext": image_ext})
        pages.append({"text": page_text, "images": images})
    doc.close()
    return pages


def to_text(input_path, output_path):
    pages = extract_with_fitz(input_path)
    if pages is None:
        print("MISSING:PyMuPDF", file=sys.stderr)
        return 2
    full = "\n\n".join([p["text"] for p in pages])
    with open(output_path, "w", encoding="utf-8") as f:
        f.write(full)
    return 0


def to_docx(input_path, output_path):
    pages = extract_with_fitz(input_path)
    if pages is None or DocxDocument is None:
        print("MISSING:docx_or_pymupdf", file=sys.stderr)
        return 2
    doc = DocxDocument()
    for i, p in enumerate(pages):
        # Add page text
        if p["text"]:
            for line in p["text"].splitlines():
                doc.add_paragraph(line)
        # Add images for that page
        for img in p["images"]:
            bio = io.BytesIO(img["bytes"])
            try:
                doc.add_picture(bio)
            except Exception:
                # fallback: write to temp file then add
                tf = output_path + f".img_{i}.{img['ext']}"
                with open(tf, "wb") as f:
                    f.write(img["bytes"])
                try:
                    doc.add_picture(tf)
                except Exception:
                    pass
                try:
                    os.remove(tf)
                except Exception:
                    pass
    doc.save(output_path)
    return 0


def to_xlsx(input_path, output_path):
    pages = extract_with_fitz(input_path)
    if pages is None or Workbook is None:
        print("MISSING:xlsx_or_pymupdf", file=sys.stderr)
        return 2
    wb = Workbook()
    # Put text in first sheet
    ws = wb.active
    ws.title = "Text"
    full_text = "\n\n".join([p["text"] for p in pages])
    ws["A1"] = full_text
    # Add images in subsequent sheets
    for i, p in enumerate(pages):
        if not p["images"]:
            continue
        sheet = wb.create_sheet(title=f"Page_{i+1}")
        col = 1
        row = 1
        for j, img in enumerate(p["images"]):
            tf = output_path + f".img_{i}_{j}.{img['ext']}"
            with open(tf, "wb") as f:
                f.write(img["bytes"])
            try:
                imgobj = XLImage(tf)
                # place image at cell A1, A10, etc.
                anchor = f"A{row}"
                sheet.add_image(imgobj, anchor)
            except Exception:
                pass
            try:
                os.remove(tf)
            except Exception:
                pass
            row += 20
    wb.save(output_path)
    return 0


def to_pptx(input_path, output_path):
    pages = extract_with_fitz(input_path)
    if pages is None or Presentation is None:
        print("MISSING:pptx_or_pymupdf", file=sys.stderr)
        return 2
    prs = Presentation()
    
    # helper: split long text into slide-sized chunks
    def split_text_chunks(text, max_chars=900):
        if not text:
            return []
        paragraphs = [pq.strip() for pq in text.split("\n\n") if pq.strip()]
        chunks = []
        cur = []
        cur_len = 0
        for p in paragraphs:
            if cur_len + len(p) + 2 > max_chars and cur:
                chunks.append("\n\n".join(cur))
                cur = [p]
                cur_len = len(p)
            else:
                cur.append(p)
                cur_len += len(p) + 2
        if cur:
            chunks.append("\n\n".join(cur))
        return chunks

    for page_index, p in enumerate(pages):
        # If page has images, create one slide per image (preserve order)
        if p["images"]:
            for img_i, img in enumerate(p["images"]):
                slide_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0]
                slide = prs.slides.add_slide(slide_layout)
                bio = io.BytesIO(img["bytes"])
                # add image filling width, keep aspect
                try:
                    slide.shapes.add_picture(bio, Inches(0.5), Inches(0.5), width=Inches(9))
                except Exception:
                    tf = output_path + f".slide_img.{page_index}_{img_i}.{img.get('ext','png')}"
                    with open(tf, "wb") as f:
                        f.write(img["bytes"])
                    try:
                        slide.shapes.add_picture(tf, Inches(0.5), Inches(0.5), width=Inches(9))
                    except Exception:
                        pass
                    try:
                        os.remove(tf)
                    except Exception:
                        pass
            # after adding image slides, also add text slides for the page text (if any)
            text_chunks = split_text_chunks(p.get("text", ""))
            for chunk in text_chunks:
                slide_layout = prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0]
                slide = prs.slides.add_slide(slide_layout)
                left = Inches(0.5)
                top = Inches(0.5)
                width = Inches(9)
                height = Inches(6)
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                tf.word_wrap = True
                paras = chunk.split("\n\n")
                for i, para in enumerate(paras):
                    ppara = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
                    ppara.text = para[:10000]
                    try:
                        ppara.font.size = Pt(18)
                    except Exception:
                        pass
        else:
            # No images: split the page text into multiple slides
            text_chunks = split_text_chunks(p.get("text", ""))
            if not text_chunks:
                # create an empty slide to preserve pagination
                slide_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
                prs.slides.add_slide(slide_layout)
            for chunk in text_chunks:
                slide_layout = prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0]
                slide = prs.slides.add_slide(slide_layout)
                left = Inches(0.5)
                top = Inches(0.5)
                width = Inches(9)
                height = Inches(6)
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                tf.word_wrap = True
                paras = chunk.split("\n\n")
                for i, para in enumerate(paras):
                    ppara = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
                    ppara.text = para[:10000]
                    try:
                        ppara.font.size = Pt(18)
                    except Exception:
                        pass
    prs.save(output_path)
    return 0


def main():
    parser = argparse.ArgumentParser(description="Convert PDF to other formats")
    parser.add_argument("input", help="Input PDF file path")
    parser.add_argument("format", help="Output format: txt|docx|xlsx|pptx")
    parser.add_argument("output", help="Output file path")
    args = parser.parse_args()

    input_path = args.input
    fmt = args.format.lower()
    output_path = args.output

    if not os.path.exists(input_path):
        print("INPUT_NOT_FOUND", file=sys.stderr)
        sys.exit(3)

    if fmt == "txt":
        rc = to_text(input_path, output_path)
    elif fmt == "docx":
        rc = to_docx(input_path, output_path)
    elif fmt == "xlsx":
        rc = to_xlsx(input_path, output_path)
    elif fmt == "pptx":
        rc = to_pptx(input_path, output_path)
    else:
        print("UNSUPPORTED_FORMAT", file=sys.stderr)
        rc = 4

    sys.exit(rc)


if __name__ == "__main__":
    main()
